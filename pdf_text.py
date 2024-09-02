from pptx import Presentation


# Function to extract text from a pptx file
def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text_content = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

    return text_content


# Define the pptx file path (make sure it's a valid path)
pptx_file = "UNIT 2 Types of errors.pptx"

# Store the extracted text in a global variable
extracted_text = extract_text_from_pptx(pptx_file)
